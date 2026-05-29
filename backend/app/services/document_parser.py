import csv
import io
from pathlib import Path
from typing import List

from langchain_text_splitters import RecursiveCharacterTextSplitter


def parse_file(file_path: str, file_type: str) -> str:
    path = Path(file_path)
    ext = file_type.lower().lstrip(".")

    if ext == "pdf":
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if ext in ("docx", "doc"):
        from docx import Document

        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext in ("pptx", "ppt"):
        from pptx import Presentation

        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

    if ext == "csv":
        with open(file_path, newline="", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            return "\n".join(", ".join(row) for row in reader)

    if ext in ("txt", "md", "markdown"):
        return path.read_text(encoding="utf-8", errors="ignore")

    return path.read_text(encoding="utf-8", errors="ignore")


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 150) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_text(text)


def extract_text_from_bytes(content: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower().lstrip(".")
    if ext == "pdf":
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if ext == "csv":
        text = content.decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(text))
        return "\n".join(", ".join(row) for row in reader)
    return content.decode("utf-8", errors="ignore")
